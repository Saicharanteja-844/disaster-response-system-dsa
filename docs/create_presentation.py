import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # Helper Functions
    # ----------------------------------------------------
    def set_dark_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 15, 25) # Dark Navy #0b0f19
        
    def add_slide_title(slide, main_text, highlight_text):
        # Vertical coral colored accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.8)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(239, 68, 68) # Coral red #ef4444
        bar.line.fill.background()
        
        # Title text box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = main_text + " "
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(241, 245, 249) # Slate 50
        
        run = p.add_run()
        run.text = highlight_text
        run.font.name = "Arial"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(6, 182, 212) # Cyan #06b6d4

    def add_card(slide, left, top, width, height, title, content_list, title_color=RGBColor(6, 182, 212), card_bg=RGBColor(21, 32, 54)):
        # Draw rounded rectangle for card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = RGBColor(38, 50, 75) # Border Slate 700
        card.line.width = Pt(1.5)
        
        # Content text box inside card
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(12)
        
        for item in content_list:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.name = "Times New Roman"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(200, 205, 215) # Light gray
            p.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Split layout with image)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide1)
    
    # Add generated rescue boat image on the right half
    image_path = r"C:\Users\india\.gemini\antigravity\brain\0d0d54cc-a1c5-4509-882d-7ff5eca40b58\rescue_boat_1784442951966.jpg"
    if os.path.exists(image_path):
        slide1.shapes.add_picture(image_path, Inches(6.666), Inches(0), Inches(6.666), Inches(7.5))
        
    # Text details on the left half
    tb = slide1.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(5.8), Inches(6.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_tag = tf.paragraphs[0]
    p_tag.text = "// NATIONAL DISASTER RESPONSE COMMAND"
    p_tag.font.name = "Arial"
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = RGBColor(200, 205, 215)
    p_tag.space_after = Pt(20)
    
    p_title = tf.add_paragraph()
    p_title.text = "RESCUE"
    p_title.font.name = "Arial"
    p_title.font.size = Pt(48)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)
    
    run_nav = p_title.add_run()
    run_nav.text = "NAV\n"
    run_nav.font.color.rgb = RGBColor(239, 68, 68) # Coral red
    
    run_sys = p_title.add_run()
    run_sys.text = "SYSTEM"
    run_sys.font.color.rgb = RGBColor(255, 255, 255)
    p_title.space_after = Pt(20)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Graph-Based Emergency Pathfinding Engine & Interactive Web Dashboard."
    p_sub.font.name = "Times New Roman"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = RGBColor(200, 205, 215)
    p_sub.space_after = Pt(40)
    
    p_dev = tf.add_paragraph()
    p_dev.text = "DEVELOPED BY:\n"
    p_dev.font.name = "Arial"
    p_dev.font.size = Pt(11)
    p_dev.font.bold = True
    p_dev.font.color.rgb = RGBColor(255, 255, 255)
    
    run_names = p_dev.add_run()
    run_names.text = "Parasa Bhavyaswi (12502256)\nPudi Sai Charan Teja (12507339)\n\n"
    run_names.font.name = "Times New Roman"
    run_names.font.size = Pt(11)
    run_names.font.bold = False
    run_names.font.color.rgb = RGBColor(200, 205, 215)
    
    p_team = tf.add_paragraph()
    p_team.text = "TEAM 9PV33:\n"
    p_team.font.name = "Arial"
    p_team.font.size = Pt(11)
    p_team.font.bold = True
    p_team.font.color.rgb = RGBColor(255, 255, 255)
    
    run_team_names = p_team.add_run()
    run_team_names.text = "Ganta Gangadhar (12623695)\nKola Midhun Kumar (12524467)\n\n"
    run_team_names.font.name = "Times New Roman"
    run_team_names.font.size = Pt(11)
    run_team_names.font.bold = False
    run_team_names.font.color.rgb = RGBColor(200, 205, 215)
    
    p_foot = tf.add_paragraph()
    p_foot.text = "LPU • DEPT. OF CSE (AI & ML) • 2026-27"
    p_foot.font.name = "Arial"
    p_foot.font.size = Pt(11)
    p_foot.font.bold = True
    p_foot.font.color.rgb = RGBColor(245, 158, 11) # Gold

    # ----------------------------------------------------
    # SLIDE 2: Project Overview
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide2)
    add_slide_title(slide2, "PROJECT", "OVERVIEW")
    
    card_width = Inches(3.7)
    card_height = Inches(4.5)
    card_y = Inches(1.8)
    card_x_starts = [Inches(0.6), Inches(4.7), Inches(8.8)]
    
    add_card(
        slide2, card_x_starts[0], card_y, card_width, card_height,
        "PROBLEM",
        [
            "Natural disasters cause severe transport blockages, cutting off critical road links.",
            "Rescue coordinators lack visibility on road passability and dynamic safety levels.",
            "No existing open-source tool integrates routing with local hospital bed and shelter capacities."
        ],
        title_color=RGBColor(239, 68, 68)
    )
    
    add_card(
        slide2, card_x_starts[1], card_y, card_width, card_height,
        "OBJECTIVES",
        [
            "Model 36 cities as weighted graphs using adjacency lists for memory efficiency.",
            "Implement and compare BFS, DFS, Dijkstra, and A* algorithms for optimal route planning.",
            "Provide interactive web visualization of shortest rescue routes and resource locations."
        ],
        title_color=RGBColor(6, 182, 212)
    )
    
    add_card(
        slide2, card_x_starts[2], card_y, card_width, card_height,
        "KEY NEEDS",
        [
            "India faces 27+ disaster types annually, requiring rapid emergency response.",
            "Manual routing and coordination fail under high-stress crisis scenarios.",
            "An integrated, DSA-powered navigation system is a national necessity for safety and efficiency."
        ],
        title_color=RGBColor(245, 158, 11)
    )

    # ----------------------------------------------------
    # SLIDE 3: Technology Stack
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide3)
    add_slide_title(slide3, "TECHNOLOGY", "STACK")
    
    card_width = Inches(5.8)
    card_height = Inches(3.6)
    card_y = Inches(1.8)
    
    add_card(
        slide3, Inches(0.6), card_y, card_width, card_height,
        "FRONTEND & VIZ",
        [
            "HTML5 / CSS3 / JavaScript (ES6+): Built a clean, responsive, dark-themed user dashboard.",
            "Leaflet.js & OpenStreetMap: Renders interactive map canvas, route polylines, and dynamic resource markers.",
            "Operations Dashboard: Unified control center for disaster toggles and routing selections."
        ]
    )
    
    add_card(
        slide3, Inches(6.9), card_y, card_width, card_height,
        "BACKEND ENGINE",
        [
            "C++ STL (Min-Heap, Priority Queue): Executes core routing and connectivity calculations.",
            "Node.js & Express REST API Layer: Standardizes communications between client and database/engine.",
            "Compiled C++ Binary Bridge: Direct stdin/stdout execution bridge (DisasterSystem.exe) for local speed."
        ],
        title_color=RGBColor(245, 158, 11)
    )
    
    # Bottom capsules
    capsules_y = Inches(5.8)
    capsule_box = slide3.shapes.add_textbox(Inches(0.6), capsules_y, Inches(12.133), Inches(0.8))
    tf = capsule_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run1 = p.add_run()
    run1.text = " DATA STORAGE: "
    run1.font.bold = True
    run1.font.name = "Arial"
    run1.font.size = Pt(11)
    run1.font.color.rgb = RGBColor(200, 205, 215)
    run1_val = p.add_run()
    run1_val.text = "CSV / JSON   |   "
    run1_val.font.bold = True
    run1_val.font.name = "Arial"
    run1_val.font.size = Pt(11)
    run1_val.font.color.rgb = RGBColor(6, 182, 212)
    
    run2 = p.add_run()
    run2.text = "DEV TOOLS: "
    run2.font.bold = True
    run2.font.name = "Arial"
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(200, 205, 215)
    run2_val = p.add_run()
    run2_val.text = "VS CODE / GIT   |   "
    run2_val.font.bold = True
    run2_val.font.name = "Arial"
    run2_val.font.size = Pt(11)
    run2_val.font.color.rgb = RGBColor(6, 182, 212)
    
    run3 = p.add_run()
    run3.text = "API TYPE: "
    run3.font.bold = True
    run3.font.name = "Arial"
    run3.font.size = Pt(11)
    run3.font.color.rgb = RGBColor(200, 205, 215)
    run3_val = p.add_run()
    run3_val.text = "RESTful"
    run3_val.font.bold = True
    run3_val.font.name = "Arial"
    run3_val.font.size = Pt(11)
    run3_val.font.color.rgb = RGBColor(6, 182, 212)

    # ----------------------------------------------------
    # SLIDE 4: System Architecture
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide4)
    add_slide_title(slide4, "SYSTEM", "ARCHITECTURE")
    
    tb = slide4.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    layers = [
        ("User Layer", "Emergency Responders & Coordinators input disaster details and request routing via the web dashboard.", RGBColor(6, 182, 212)),
        ("API Layer", "Node.js & Express server receives client inputs, parses JSON data, and handles command-line arguments.", RGBColor(239, 68, 68)),
        ("C++ Pathfinding Engine", "Core algorithmic module utilizing optimized C++ structures to compute shortest path in sub-milliseconds.", RGBColor(245, 158, 11)),
        ("Data Layer", "Loads and updates CSV/JSON databases for road maps, hospitals (bed counts), shelters, and active response teams.", RGBColor(16, 185, 129))
    ]
    
    for idx, (title, desc, color) in enumerate(layers):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.text = f"{title}: "
        p.font.name = "Arial"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(4)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Times New Roman"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RGBColor(200, 205, 215)
        p_desc.space_after = Pt(14)
        
    # Flow Diagram blocks
    block_x = Inches(7.5)
    block_width = Inches(5.2)
    block_height = Inches(0.9)
    block_y_starts = [Inches(1.8), Inches(3.1), Inches(4.4), Inches(5.7)]
    
    block_names = [
        ("USER LAYER", "Web Dashboard • Emergency Responders • Coordinators", RGBColor(6, 182, 212)),
        ("API LAYER", "Node.js + Express • REST Endpoints • Request Routing", RGBColor(239, 68, 68)),
        ("C++ PATHFINDING ENGINE", "BFS • DFS • Dijkstra • A* • 36-Node City Graph", RGBColor(245, 158, 11)),
        ("DATA LAYER", "Hospitals • Shelters • Rescue Teams • Live Locations", RGBColor(16, 185, 129))
    ]
    
    for idx, (b_name, b_sub, b_color) in enumerate(block_names):
        rect = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, block_x, block_y_starts[idx], block_width, block_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(21, 32, 54)
        rect.line.color.rgb = b_color
        rect.line.width = Pt(1.5)
        
        btf = rect.text_frame
        btf.word_wrap = True
        btf.margin_top = btf.margin_bottom = Inches(0.08)
        p1 = btf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = b_name
        p1.font.name = "Arial"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = b_color
        
        p2 = btf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = b_sub
        p2.font.name = "Arial"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(200, 205, 215)
        
        if idx < 3:
            arrow = slide4.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                block_x + block_width/2 - Inches(0.15),
                block_y_starts[idx] + block_height + Inches(0.05),
                Inches(0.3), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(6, 182, 212)
            arrow.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 5: Pathfinding Algorithms
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide5)
    add_slide_title(slide5, "PATHFINDING", "ALGORITHMS")
    
    card_width = Inches(3.7)
    card_height = Inches(4.5)
    card_y = Inches(1.8)
    card_x_starts = [Inches(0.6), Inches(4.7), Inches(8.8)]
    
    add_card(
        slide5, card_x_starts[0], card_y, card_width, card_height,
        "BFS / DFS",
        [
            "Used for Connectivity Analysis of the road network.",
            "BFS explores level-by-level (using a Queue) to find shortest path in terms of hops.",
            "DFS explores recursively (using a Stack) to check overall graph connectivity.",
            "COMPLEXITY: O(V + E)"
        ],
        title_color=RGBColor(239, 68, 68)
    )
    
    add_card(
        slide5, card_x_starts[1], card_y, card_width, card_height,
        "DIJKSTRA",
        [
            "The Shortest Path specialist for weighted graphs.",
            "Computes safety-aware paths utilizing a Min-Heap priority queue to fetch the lowest cost node.",
            "Safety Weight Formula:\nEdge Cost = Distance * (1.0 + (Danger Level - 1) * 0.3)",
            "COMPLEXITY: O(E log V)"
        ],
        title_color=RGBColor(6, 182, 212)
    )
    
    add_card(
        slide5, card_x_starts[2], card_y, card_width, card_height,
        "A* SEARCH",
        [
            "Goal-directed pathfinding utilizing heuristic search f(n) = g(n) + h(n).",
            "Heuristic h(n) uses the Haversine formula to compute great-circle distance using city coordinates.",
            "Reduces search space and expanded nodes by up to 56% compared to Dijkstra.",
            "COMPLEXITY: O(E log V)"
        ],
        title_color=RGBColor(245, 158, 11)
    )

    # ----------------------------------------------------
    # SLIDE 6: Data Structures
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide6)
    add_slide_title(slide6, "DATA", "STRUCTURES")
    
    card_width = Inches(5.8)
    card_height = Inches(3.6)
    card_y = Inches(1.8)
    
    add_card(
        slide6, Inches(0.6), card_y, card_width, card_height,
        "GRAPH REPRESENTATION",
        [
            "Adjacency List: Models the transportation network of 41 cities and 69 connecting highways.",
            "Highly memory-efficient O(V + E) storage structure, ideal for sparse road systems.",
            "Dynamically scalable: Allows instant modification of edge weights and danger ratings.",
            "Alphabetical sorting of adjacency lists ensures determinism in traversals."
        ]
    )
    
    add_card(
        slide6, Inches(6.9), card_y, card_width, card_height,
        "MEMORY MANAGEMENT",
        [
            "Vector: Used for dynamic path storage, offering contiguous memory and O(1) random access.",
            "Priority Queue: Backed by binary min-heap for O(log V) retrieval of minimum-distance node.",
            "Hash Maps (unordered_map): Provides O(1) average lookup for city coordinates and metadata.",
            "Stack & Queue: Standard DFS recursion stack and BFS level-order queues."
        ],
        title_color=RGBColor(245, 158, 11)
    )
    
    # Formula box at bottom
    formula_y = Inches(5.8)
    formula_box = slide6.shapes.add_textbox(Inches(0.6), formula_y, Inches(12.133), Inches(0.8))
    ftf = formula_box.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fp.text = "Distance Formula (u, v) = √[(x_v - x_u)² + (y_v - y_u)²]"
    fp.font.name = "Arial"
    fp.font.size = Pt(18)
    fp.font.bold = True
    fp.font.color.rgb = RGBColor(6, 182, 212)

    # ----------------------------------------------------
    # SLIDE 7: Core Modules
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide7)
    add_slide_title(slide7, "CORE", "MODULES")
    
    grid_width = Inches(5.8)
    grid_height = Inches(2.2)
    y_starts = [Inches(1.8), Inches(4.3)]
    x_starts = [Inches(0.6), Inches(6.9)]
    
    modules = [
        ("Operations Dashboard", [
            "Live statistics loaded from JSON datasets.",
            "Provides dynamic weather info, routing toggles, and algorithm selection panels.",
            "Visualizes active disaster zones and affected highway coordinates."
        ], RGBColor(6, 182, 212)),
        ("Route Planning", [
            "Selects source and destination cities interactively.",
            "Executes C++ engine pathfinders locally or falls back to serverless JS router.",
            "Supports dynamic road blockage simulation by setting safety danger scores."
        ], RGBColor(239, 68, 68)),
        ("Emergency Resources", [
            "Scans for nearby hospitals and lists active phone lines and bed counts.",
            "Identifies designated emergency shelters with capacities and occupancy levels.",
            "Matches and details local search-and-rescue teams (SDRF) for immediate dispatch."
        ], RGBColor(245, 158, 11)),
        ("Dataset Management", [
            "Manages files in comma-separated values (CSV) formats.",
            "Reads and writes new data, updates road maps, and auto-rebuilds the graph.",
            "Integrates hierarchical JSON datasets for active disaster alerts."
        ], RGBColor(16, 185, 129))
    ]
    
    for idx, (title, contents, color) in enumerate(modules):
        x = x_starts[idx % 2]
        y = y_starts[idx // 2]
        add_card(slide7, x, y, grid_width, grid_height, title, contents, title_color=color)

    # ----------------------------------------------------
    # SLIDE 8: Execution Results (With generated mock dashboard)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide8)
    add_slide_title(slide8, "EXECUTION", "RESULTS")
    
    # Left side: image
    img_path = r"C:\Users\india\.gemini\antigravity\brain\0d0d54cc-a1c5-4509-882d-7ff5eca40b58\mock_dashboard_1784442992711.jpg"
    if os.path.exists(img_path):
        slide8.shapes.add_picture(img_path, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
        
    # Right side: results
    tb = slide8.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p1 = tf.paragraphs[0]
    p1.text = "Optimal Routing: "
    p1.font.name = "Arial"
    p1.font.size = Pt(17)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(6, 182, 212)
    p1.space_after = Pt(3)
    
    p1_desc = tf.add_paragraph()
    p1_desc.text = "Successfully calculated Jammu to Bengaluru route (2790 km) in 3350 mins. The engine bypassed active flood zones dynamically."
    p1_desc.font.name = "Times New Roman"
    p1_desc.font.size = Pt(13)
    p1_desc.font.color.rgb = RGBColor(200, 205, 215)
    p1_desc.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "Latency: "
    p2.font.name = "Arial"
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(239, 68, 68)
    p2.space_after = Pt(3)
    
    p2_desc = tf.add_paragraph()
    p2_desc.text = "Core C++ engine average pathfinding latency: 0.04 ms. Serverless cloud fallback execution latency: under 12 ms."
    p2_desc.font.name = "Times New Roman"
    p2_desc.font.size = Pt(13)
    p2_desc.font.color.rgb = RGBColor(200, 205, 215)
    p2_desc.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "Adaptability: "
    p3.font.name = "Arial"
    p3.font.size = Pt(17)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(245, 158, 11)
    p3.space_after = Pt(3)
    
    p3_desc = tf.add_paragraph()
    p3_desc.text = "System dynamically recalculates alternative paths when specific highways are blocked (e.g., Patna-Kolkata corridor)."
    p3_desc.font.name = "Times New Roman"
    p3_desc.font.size = Pt(13)
    p3_desc.font.color.rgb = RGBColor(200, 205, 215)
    p3_desc.space_after = Pt(26)
    
    p_cap = tf.add_paragraph()
    p_cap.text = "CITIES: "
    p_cap.font.name = "Arial"
    p_cap.font.size = Pt(13)
    p_cap.font.bold = True
    p_cap.font.color.rgb = RGBColor(255, 255, 255)
    
    run_city = p_cap.add_run()
    run_city.text = "36   |   "
    run_city.font.color.rgb = RGBColor(6, 182, 212)
    
    run_acc_lbl = p_cap.add_run()
    run_acc_lbl.text = "ACCURACY: "
    run_acc_lbl.font.bold = True
    
    run_acc = p_cap.add_run()
    run_acc.text = "100%"
    run_acc.font.color.rgb = RGBColor(16, 185, 129)

    # ----------------------------------------------------
    # SLIDE 9: Performance Metrics Table
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide9)
    add_slide_title(slide9, "PERFORMANCE", "METRICS")
    
    # Add Table
    left = Inches(1.0)
    top = Inches(1.8)
    width = Inches(11.333)
    height = Inches(4.5)
    
    table_shape = slide9.shapes.add_table(5, 5, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(2.2) # Algorithm
    table.columns[1].width = Inches(2.2) # Complexity
    table.columns[2].width = Inches(2.2) # Exec Time
    table.columns[3].width = Inches(2.2) # Optimal Path
    table.columns[4].width = Inches(2.533) # Visited Nodes
    
    headers = ["ALGORITHM", "COMPLEXITY", "EXEC TIME", "OPTIMAL PATH", "VISITED NODES"]
    data = [
        ["BFS", "O(V + E)", "48 ms", "NO", "30 Nodes"],
        ["DFS", "O(V + E)", "31 ms", "NO", "27 Nodes"],
        ["Dijkstra", "O(E log V)", "12 ms", "YES", "25 Nodes"],
        ["A* Search", "O(E log V)", "8 ms", "YES", "19 Nodes"]
    ]
    
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(21, 32, 54)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(6, 182, 212)
        
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            else:
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
                
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Times New Roman"
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(200, 205, 215)
            
            if text == "YES":
                p.font.color.rgb = RGBColor(16, 185, 129)
                p.font.bold = True
            elif text == "NO":
                p.font.color.rgb = RGBColor(239, 68, 68)
                p.font.bold = True

    # ----------------------------------------------------
    # SLIDE 10: System Strengths
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_dark_background(slide10)
    add_slide_title(slide10, "SYSTEM", "STRENGTHS")
    
    card_width = Inches(3.7)
    card_height = Inches(2.2)
    y_starts = [Inches(1.8), Inches(4.3)]
    x_starts = [Inches(0.6), Inches(4.7), Inches(8.8)]
    
    strengths = [
        ("Instant Pathfinding", "Dijkstra/A* execute in milliseconds across 36 nodes for emergency response.", RGBColor(6, 182, 212)),
        ("Dynamic Obstacles", "Real-time road blockage management dynamically bypasses danger zones.", RGBColor(239, 68, 68)),
        ("Resource Mapping", "Locates hospitals (beds available) and emergency shelters along the route.", RGBColor(245, 158, 11)),
        ("Modular Design", "High-performance C++ backend bridged to Node.js handles scalability and cloud fallback.", RGBColor(16, 185, 129)),
        ("Interactive UI", "Real-time route polyline and resource markers rendered dynamically via Leaflet.js.", RGBColor(167, 139, 250)),
        ("Validated", "System passed all test cases (TC-01 to TC-09) with a 100% success rate in validation.", RGBColor(251, 146, 60))
    ]
    
    for idx, (title, desc, color) in enumerate(strengths):
        x = x_starts[idx % 3]
        y = y_starts[idx // 3]
        add_card(slide10, x, y, card_width, card_height, title, [desc], title_color=color)

    # Save presentation
    output_path = r"C:\Users\india\.gemini\antigravity\scratch\project_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

    # Copy to .ppt so user has both formats
    output_path_ppt = r"C:\Users\india\.gemini\antigravity\scratch\project_presentation.ppt"
    # We will overwrite the old HTML-pretending project_presentation.ppt
    try:
        import shutil
        shutil.copy(output_path, output_path_ppt)
        print(f"Copied presentation to: {output_path_ppt}")
    except Exception as e:
        print(f"Failed to copy to .ppt: {e}")

if __name__ == "__main__":
    main()
